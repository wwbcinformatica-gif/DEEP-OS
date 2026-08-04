"""
Document Reader — Le PDF, DOCX, XLSX, CSV, XML, TXT, PPTX e outros formatos.
"""
import os
import csv
import json
import xml.etree.ElementTree as ET


async def tool_read_document(path: str, root: str = "") -> dict:
    """Le o conteudo de documentos em varios formatos."""
    base = root or os.getcwd()
    full_path = os.path.normpath(os.path.join(base, path)) if not os.path.isabs(path) else path

    if not os.path.exists(full_path):
        return {"error": f"Arquivo nao encontrado: {full_path}"}

    ext = os.path.splitext(full_path)[1].lower()

    try:
        if ext == ".pdf":
            return _read_pdf(full_path)
        elif ext in (".docx", ".doc"):
            return _read_docx(full_path)
        elif ext in (".xlsx", ".xls"):
            return _read_xlsx(full_path)
        elif ext == ".csv":
            return _read_csv(full_path)
        elif ext in (".xml", ".svg", ".xhtml", ".html", ".htm"):
            return _read_xml(full_path)
        elif ext in (".pptx", ".ppt"):
            return _read_pptx(full_path)
        elif ext in (".txt", ".md", ".py", ".js", ".ts", ".tsx", ".jsx",
                      ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg",
                      ".log", ".env", ".sh", ".bat", ".ps1", ".css", ".scss",
                      ".sql", ".r", ".java", ".c", ".cpp", ".h", ".hpp",
                      ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".scala"):
            return _read_text(full_path)
        else:
            return _read_text(full_path)
    except Exception as e:
        return {"error": f"Erro ao ler {ext}: {str(e)}"}


def _read_pdf(path: str) -> dict:
    """Le PDF com pdfplumber (melhor extração de texto)."""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = []
            total_chars = 0
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                # Limpar quebras de linha desnecessarias
                text = text.replace('\n', ' ').strip()
                # Reformatar em paragrafos
                text = '  \n'.join([p.strip() for p in text.split('  ') if p.strip()])
                total_chars += len(text)
                pages.append({"page": i + 1, "text": text})
            full_text = "\n\n".join([p["text"] for p in pages])
            return {
                "format": "PDF",
                "file": os.path.basename(path),
                "pages": len(pdf.pages),
                "total_chars": total_chars,
                "content": full_text[:50000],
                "truncated": total_chars > 50000,
            }
    except ImportError:
        # Fallback para PyPDF2
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        pages = []
        total_chars = 0
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.replace('\n', ' ').strip()
            total_chars += len(text)
            pages.append({"page": i + 1, "text": text})
        full_text = "\n\n".join([p["text"] for p in pages])
        return {
            "format": "PDF",
            "file": os.path.basename(path),
            "pages": len(reader.pages),
            "total_chars": total_chars,
            "content": full_text[:50000],
            "truncated": total_chars > 50000,
        }


def _read_docx(path: str) -> dict:
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Extrair tabelas
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append(cells)
        tables.append(rows)
    full_text = "\n".join(paragraphs)
    return {
        "format": "DOCX",
        "file": os.path.basename(path),
        "paragraphs": len(paragraphs),
        "tables": len(tables),
        "content": full_text[:50000],
        "tables_data": tables[:10] if tables else [],
        "truncated": len(full_text) > 50000,
    }


def _read_xlsx(path: str) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    sheets = {}
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(max_row=200, values_only=True):
            rows.append([str(c) if c is not None else "" for c in row])
        sheets[name] = rows
    wb.close()
    total_rows = sum(len(r) for r in sheets.values())
    return {
        "format": "XLSX",
        "file": os.path.basename(path),
        "sheet_names": list(sheets.keys()),
        "total_rows": total_rows,
        "sheets": sheets,
    }


def _read_csv(path: str) -> dict:
    rows = []
    encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]
    content = None
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc) as f:
                content = f.read()
                break
        except (UnicodeDecodeError, UnicodeError):
            continue
    if content is None:
        return {"error": "Nao foi possivel decodificar o CSV"}

    reader = csv.reader(content.splitlines())
    for i, row in enumerate(reader):
        if i >= 500:
            break
        rows.append(row)

    return {
        "format": "CSV",
        "file": os.path.basename(path),
        "total_rows": len(rows),
        "headers": rows[0] if rows else [],
        "content": rows[:500],
    }


def _read_xml(path: str) -> dict:
    import chardet
    with open(path, "rb") as f:
        raw = f.read()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8")
    text = raw.decode(encoding, errors="replace")

    try:
        root = ET.fromstring(text)
        elements = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                elements.append({"tag": elem.tag, "text": elem.text.strip()[:500]})
        return {
            "format": "XML",
            "file": os.path.basename(path),
            "root_tag": root.tag,
            "elements_count": len(elements),
            "content": text[:50000],
            "elements": elements[:200],
            "truncated": len(text) > 50000,
        }
    except ET.ParseError:
        return {
            "format": "XML (raw)",
            "file": os.path.basename(path),
            "content": text[:50000],
            "truncated": len(text) > 50000,
        }


def _read_pptx(path: str) -> dict:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append({"slide": i + 1, "text": "\n".join(texts)})
    full_text = "\n\n".join([s["text"] for s in slides])
    return {
        "format": "PPTX",
        "file": os.path.basename(path),
        "slides": len(slides),
        "content": full_text[:50000],
        "truncated": len(full_text) > 50000,
    }


def _read_text(path: str) -> dict:
    import chardet
    with open(path, "rb") as f:
        raw = f.read()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding", "utf-8")
    text = raw.decode(encoding, errors="replace")
    ext = os.path.splitext(path)[1].lower()
    return {
        "format": ext.upper().replace(".", ""),
        "file": os.path.basename(path),
        "encoding": encoding,
        "lines": text.count("\n") + 1,
        "total_chars": len(text),
        "content": text[:50000],
        "truncated": len(text) > 50000,
    }
